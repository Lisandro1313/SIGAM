"""Generate SIGAM project presentation (PPTX)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Brand colors ──
PRIMARY = RGBColor(0x1B, 0x5E, 0x20)      # dark green
ACCENT = RGBColor(0x2E, 0x7D, 0x32)       # green
LIGHT_BG = RGBColor(0xE8, 0xF5, 0xE9)     # light green bg
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x61, 0x61, 0x61)
BLUE = RGBColor(0x15, 0x65, 0xC0)
ORANGE = RGBColor(0xE6, 0x51, 0x00)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = spacing
        p.level = 0
    return txBox


def slide_header(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.15), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.18), Inches(10), Inches(0.7),
                 title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.4),
                     subtitle, font_size=16, color=RGBColor(0xA5, 0xD6, 0xA7))


def card(slide, left, top, width, height, title, items, icon="", title_color=ACCENT):
    # card bg
    shape = add_rect(slide, left, top, width, height, WHITE)
    shape.shadow.inherit = False
    # title
    t = f"{icon}  {title}" if icon else title
    add_text_box(slide, left + Inches(0.25), top + Inches(0.15), width - Inches(0.5), Inches(0.45),
                 t, font_size=18, color=title_color, bold=True)
    # items
    add_bullet_list(slide, left + Inches(0.25), top + Inches(0.6), width - Inches(0.5),
                    height - Inches(0.8), items, font_size=13, color=GRAY, spacing=Pt(4))


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, PRIMARY)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height,
         RGBColor(0x14, 0x4D, 0x17))

# Decorative accent bar
add_rect(slide, Inches(1.5), Inches(2.4), Inches(1.2), Inches(0.06), WHITE)

add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(1.2),
             "SIGAM", font_size=72, color=WHITE, bold=True)
add_text_box(slide, Inches(1.5), Inches(3.7), Inches(10), Inches(0.8),
             "Sistema Integral de Gestión Alimentaria Municipal",
             font_size=28, color=RGBColor(0xA5, 0xD6, 0xA7))
add_text_box(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.5),
             "Municipalidad de La Plata  —  Secretaría de Desarrollo Social",
             font_size=18, color=RGBColor(0x81, 0xC7, 0x84))

add_rect(slide, Inches(1.5), Inches(5.5), Inches(1.2), Inches(0.06), WHITE)
add_text_box(slide, Inches(1.5), Inches(5.7), Inches(10), Inches(0.4),
             "Presentación de Proyecto  ·  2026",
             font_size=16, color=RGBColor(0x81, 0xC7, 0x84))

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — Problema
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_header(slide, "El Problema", "¿Por qué se necesita SIGAM?")

problems = [
    ("Planillas Excel manuales", "Errores de carga, datos duplicados, sin control de versiones"),
    ("Sin trazabilidad", "No se puede verificar quién recibió qué, cuándo ni dónde"),
    ("Stock sin control", "Mercadería vencida, faltantes no detectados, transferencias informales"),
    ("Casos urgentes sin proceso", "Pedidos por WhatsApp sin registro ni seguimiento"),
    ("Reportes manuales", "ANEXO VI provincial armado a mano, propenso a errores"),
    ("Cero visibilidad en tiempo real", "Sin dashboards ni alertas automáticas"),
]

for i, (title, desc) in enumerate(problems):
    col = i % 3
    row = i // 3
    left = Inches(0.6 + col * 4.1)
    top = Inches(1.6 + row * 2.7)
    shape = add_rect(slide, left, top, Inches(3.8), Inches(2.3), WHITE)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.25), Inches(3.2), Inches(0.5),
                 f"✗  {title}", font_size=18, color=ORANGE, bold=True)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.85), Inches(3.2), Inches(1.2),
                 desc, font_size=14, color=GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — Solución
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_header(slide, "La Solución: SIGAM", "Plataforma web integral de gestión alimentaria")

solutions = [
    ("Gestión Centralizada", [
        "Base de datos única y en la nube",
        "Acceso multi-usuario simultáneo",
        "Datos consistentes y sin duplicados",
        "Búsqueda global por DNI (Ctrl+K)",
    ]),
    ("Trazabilidad Completa", [
        "Cada entrega documentada con foto de firma",
        "Historial completo de movimientos de stock",
        "Auditoría automática de toda acción",
        "PDF de remitos descargables",
    ]),
    ("Automatización", [
        "Cronograma mensual autogenerado",
        "Alertas de vencimiento de mercadería",
        "Notificaciones en tiempo real (SSE)",
        "Generación masiva de remitos",
    ]),
]

for i, (title, items) in enumerate(solutions):
    left = Inches(0.6 + i * 4.1)
    card(slide, left, Inches(1.6), Inches(3.8), Inches(4.8),
         title, items, icon="✓", title_color=ACCENT)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Módulos Principales
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_header(slide, "Módulos del Sistema")

modules = [
    ("Beneficiarios", ["Alta/baja/modificación", "Documentación adjunta", "Detección duplicados DNI", "Historial de entregas"]),
    ("Remitos", ["Ciclo completo de entrega", "Generación de PDF", "Envío por email/WhatsApp", "Firma fotográfica"]),
    ("Stock & Depósitos", ["2 depósitos independientes", "Control de lotes y vencimiento", "Transferencias trazables", "Alertas automáticas"]),
    ("Cronograma", ["Generación mensual automática", "Vista semanal", "Ajuste de fechas", "Generación masiva de remitos"]),
    ("Casos Particulares", ["Workflow de aprobación", "Prioridad: Normal/Alta/Urgente", "Cruce automático por DNI", "Documentación de campo"]),
    ("Reportes & ANEXO VI", ["Dashboard en tiempo real", "Gráficos de distribución", "Ranking de artículos", "Exportación provincial oficial"]),
    ("Mapa Interactivo", ["Geolocalización de beneficiarios", "Filtros por localidad", "Popups informativos", "Exportación por zona"]),
    ("Auditoría", ["Log automático de toda acción", "Filtros por período/usuario", "Exportación a Excel", "Trazabilidad completa"]),
]

for i, (title, items) in enumerate(modules):
    col = i % 4
    row = i // 4
    left = Inches(0.4 + col * 3.15)
    top = Inches(1.5 + row * 2.85)
    card(slide, left, top, Inches(2.95), Inches(2.6), title, items, title_color=ACCENT)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Flujo de Entrega
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_header(slide, "Flujo de Entrega", "Ciclo de vida completo de un remito")

steps = [
    ("1", "Programar", "Cronograma mensual\nautogenerado"),
    ("2", "Crear Remito", "Manual o masivo\ndesde cronograma"),
    ("3", "Confirmar", "Stock descontado\nautomáticamente"),
    ("4", "Notificar", "Email y WhatsApp\nal beneficiario"),
    ("5", "Entregar", "Depósito registra\ncon foto de firma"),
    ("6", "Completar", "Registro en auditoría\ny reportes"),
]

for i, (num, title, desc) in enumerate(steps):
    left = Inches(0.5 + i * 2.1)
    top = Inches(2.5)
    # Circle with number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.6), top, Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Arrow (except last)
    if i < len(steps) - 1:
        add_text_box(slide, left + Inches(1.4), top + Inches(0.1), Inches(0.6), Inches(0.5),
                     "→", font_size=28, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, left + Inches(0.1), top + Inches(0.9), Inches(1.7), Inches(0.4),
                 title, font_size=16, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.05), top + Inches(1.35), Inches(1.8), Inches(0.8),
                 desc, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Roles y Permisos
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_header(slide, "Roles y Permisos", "Control de acceso granular por rol y secretaría")

roles = [
    ("ADMIN", "Acceso completo al sistema. Gestión de usuarios, programas y configuración global."),
    ("OPERADOR DE PROGRAMA", "Gestión de beneficiarios, remitos, cronograma y casos. Secretaría de Política Alimentaria."),
    ("LOGÍSTICA / DEPÓSITO", "Control de stock, recepción de mercadería, confirmación de entregas en depósito."),
    ("ASISTENCIA CRÍTICA", "Mismas funciones que Operador pero aislado en datos de Asistencia Crítica (depósito CITA)."),
    ("TRABAJADORA SOCIAL", "Creación de casos particulares desde el campo. Seguimiento de sus propios casos."),
    ("VISOR", "Acceso de solo lectura a dashboard, beneficiarios y reportes."),
]

for i, (role, desc) in enumerate(roles):
    col = i % 3
    row = i // 3
    left = Inches(0.6 + col * 4.1)
    top = Inches(1.6 + row * 2.7)
    shape = add_rect(slide, left, top, Inches(3.8), Inches(2.3), WHITE)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.25), Inches(3.2), Inches(0.5),
                 role, font_size=17, color=PRIMARY, bold=True)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.85), Inches(3.2), Inches(1.2),
                 desc, font_size=14, color=GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 — Arquitectura Técnica
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_header(slide, "Arquitectura Técnica", "Stack moderno, escalable y seguro")

# Frontend card
card(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(5.2),
     "Frontend", [
         "React 18 + TypeScript",
         "Material UI (MUI v5)",
         "Zustand (estado global)",
         "React Router v6",
         "Recharts (gráficos)",
         "Leaflet (mapas)",
         "PWA — instalable como app",
         "Vite (build ultrarrápido)",
     ], title_color=BLUE)

# Backend card
card(slide, Inches(4.8), Inches(1.5), Inches(3.8), Inches(5.2),
     "Backend", [
         "NestJS 10 + TypeScript",
         "Prisma ORM",
         "JWT + Passport.js",
         "16 módulos independientes",
         "Swagger / OpenAPI docs",
         "SSE (notificaciones real-time)",
         "PDFKit (generación de remitos)",
         "Rate limiting (60 req/min)",
     ], title_color=BLUE)

# Infrastructure card
card(slide, Inches(9.1), Inches(1.5), Inches(3.8), Inches(5.2),
     "Infraestructura", [
         "PostgreSQL 14+ (Supabase)",
         "~20 tablas relacionales",
         "Supabase Storage (archivos)",
         "Gmail SMTP (emails)",
         "Backups automáticos en nube",
         "HTTPS + CORS configurado",
         "Health checks (/health)",
         "Auditoría global automática",
     ], title_color=BLUE)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 — Seguridad y Compliance
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_header(slide, "Seguridad y Cumplimiento")

sec_items = [
    ("Autenticación JWT", "Tokens seguros con expiración. Contraseñas hasheadas con bcrypt. Sesiones stateless."),
    ("Control de Acceso (RBAC)", "Cada endpoint protegido por rol. Datos aislados por secretaría. Guards en todas las rutas."),
    ("Auditoría Completa", "Interceptor global registra automáticamente toda acción: usuario, fecha, datos modificados."),
    ("Rate Limiting", "Protección contra abuso: 60 peticiones por minuto por IP. Prevención de ataques de fuerza bruta."),
    ("Datos en la Nube", "PostgreSQL en Supabase con backups automáticos. Archivos en storage seguro con políticas de acceso."),
    ("Validación Estricta", "Todos los datos validados en backend con class-validator. Sanitización de inputs."),
]

for i, (title, desc) in enumerate(sec_items):
    col = i % 3
    row = i // 3
    left = Inches(0.6 + col * 4.1)
    top = Inches(1.6 + row * 2.7)
    shape = add_rect(slide, left, top, Inches(3.8), Inches(2.3), WHITE)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.25), Inches(3.2), Inches(0.5),
                 f"🔒  {title}", font_size=17, color=PRIMARY, bold=True)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.85), Inches(3.2), Inches(1.2),
                 desc, font_size=14, color=GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 — Beneficios Clave
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_header(slide, "Beneficios Clave")

benefits = [
    ("Eficiencia Operativa", "Automatización de cronogramas, remitos masivos y alertas. Menos trabajo manual, menos errores."),
    ("Transparencia Total", "Cada entrega documentada. Auditoría completa. Reportes exportables para rendición de cuentas."),
    ("Acceso Desde Cualquier Lugar", "Aplicación web + PWA. Trabajadoras sociales cargan casos desde el campo con el celular."),
    ("Toma de Decisiones", "Dashboard en tiempo real. Gráficos de distribución. Mapa interactivo de cobertura."),
    ("Cumplimiento Provincial", "ANEXO VI generado automáticamente. Datos listos para auditoría externa."),
    ("Escalabilidad", "Arquitectura modular preparada para incorporar nuevas secretarías y programas."),
]

for i, (title, desc) in enumerate(benefits):
    col = i % 3
    row = i // 3
    left = Inches(0.6 + col * 4.1)
    top = Inches(1.6 + row * 2.7)
    shape = add_rect(slide, left, top, Inches(3.8), Inches(2.3), WHITE)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.25), Inches(3.2), Inches(0.5),
                 f"✦  {title}", font_size=17, color=ACCENT, bold=True)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.85), Inches(3.2), Inches(1.2),
                 desc, font_size=14, color=GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 — Roadmap
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
slide_header(slide, "Próximos Pasos", "Evolución planificada del sistema")

phases = [
    ("Fase 1 — Consolidación", ACCENT, [
        "Optimización de rendimiento",
        "Mejoras de UX según feedback",
        "Estabilización de reportes",
    ]),
    ("Fase 2 — Integraciones", BLUE, [
        "WhatsApp Business API (Twilio/Meta)",
        "Notificaciones automáticas de retiro",
        "Informes por email programados",
    ]),
    ("Fase 3 — Expansión", ORANGE, [
        "Nuevas secretarías (CITA, Niñez)",
        "Datos cruzados por beneficiario",
        "Portal de autogestión",
    ]),
    ("Fase 4 — Inteligencia", RGBColor(0x6A, 0x1B, 0x9A), [
        "Predicción de demanda",
        "Optimización de rutas de entrega",
        "Análisis de impacto social",
    ]),
]

for i, (title, color, items) in enumerate(phases):
    left = Inches(0.4 + i * 3.2)
    top = Inches(1.6)
    # Phase header bar
    add_rect(slide, left, top, Inches(2.95), Inches(0.6), color)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.1), Inches(2.65), Inches(0.4),
                 title, font_size=15, color=WHITE, bold=True)
    # Items
    card_shape = add_rect(slide, left, top + Inches(0.6), Inches(2.95), Inches(3.5), WHITE)
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.85), Inches(2.55), Inches(3.0),
                    [f"•  {item}" for item in items], font_size=14, color=GRAY, spacing=Pt(8))

# ════════════════════════════════════════════════════════════════
# SLIDE 11 — Cierre
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0x14, 0x4D, 0x17))

add_rect(slide, Inches(3), Inches(2.2), Inches(7.3), Inches(0.06), WHITE)

add_text_box(slide, Inches(3), Inches(2.5), Inches(7.3), Inches(1),
             "SIGAM", font_size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3), Inches(3.5), Inches(7.3), Inches(0.6),
             "Gestión alimentaria transparente, eficiente y trazable",
             font_size=22, color=RGBColor(0xA5, 0xD6, 0xA7), alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(3), Inches(4.5), Inches(7.3), Inches(0.06), WHITE)

add_text_box(slide, Inches(3), Inches(4.8), Inches(7.3), Inches(0.5),
             "Municipalidad de La Plata  ·  Secretaría de Desarrollo Social",
             font_size=16, color=RGBColor(0x81, 0xC7, 0x84), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(3), Inches(5.6), Inches(7.3), Inches(0.5),
             "¿Preguntas?",
             font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ── Save ──
output_path = r"C:\Users\Usuario\OneDrive\Escritorio\SIGAM\SIGAM_Presentacion.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
